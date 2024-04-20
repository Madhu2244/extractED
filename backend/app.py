from flask import Flask, request, jsonify
from flask_cors import CORS
from pptx import Presentation
import os
import nltk
from nltk.tokenize import sent_tokenize
import re
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
import google.generativeai as genai
from langchain_openai import OpenAI
from langchain.text_splitter import CharacterTextSplitter
from langchain.docstore.document import Document
from langchain.chains.summarize import load_summarize_chain

# $env:API_KEY = 'AIzaSyDtdYPjd5cL3kE6daHpbbkgPS20OknKJDw'

# Initialize Flask App
app = Flask(__name__)
CORS(app)
nltk.download('punkt')
model = SentenceTransformer('all-MiniLM-L6-v2')
genai.configure(api_key=os.environ["API_KEY"])
gemini_model = genai.GenerativeModel('gemini-pro')


# Function to extract text
def extract_text_from_presentation(ppt_file):
    prs = Presentation(ppt_file)
    texts = ''

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts += shape.text

    return texts

# Function to split text into sentences
def extract_sentences_from_text(texts):
    texts = re.sub(r'’', "'", texts)
    sentences = sent_tokenize(texts)
    return sentences

# Combine similar sentences into a graph | dict()
def create_graph(sentences):
    graph = {sentence: [] for sentence in sentences}

    embed = model.encode(sentences)
    similarity_matrix = cosine_similarity(embed)

    for i in range(len(sentences)):
        for j in range(i + 1, len(sentences)):
            if i != j and similarity_matrix[i][j] > 0.4:
                graph[sentences[i]].append(sentences[j])
                graph[sentences[j]].append(sentences[i])

    return graph

def create_headers(graph):
    was_seen = set()
    headers = {}
    for key, sentences in graph.items():
        if key not in was_seen:
            was_seen.add(key)
            was_seen.update(sentences)

            similar_sentences = list(set([key] + sentences))
            input_text = " ".join(similar_sentences)
            
            header = (gemini_model.generate_content(f'Summarize the following sentences into a concise 3 to 4 word header for note-taking: {input_text}'))
            header_text = header.text.strip().replace('**', '')
            
            headers[header_text] = similar_sentences
    return headers

# Function to summarize clusters of sentences associated with headers
def summarize_clusters(headers):
    llm = OpenAI(temperature=0, max_tokens=600)  
    text_splitter = CharacterTextSplitter()  
    summaries = {}  

    for header, sentences in headers.items():
        input_text = " ".join(sentences)
        
        texts = text_splitter.split_text(input_text)
        docs = [Document(page_content=t) for t in texts]  
        
        chain = load_summarize_chain(llm, chain_type="map_reduce")
        
        summary_output = chain.invoke(docs)
        clean_summary = summary_output['output_text'].replace('\n', ' ').strip()
        summaries[header] = clean_summary

    return summaries

def generate_notes_from_summaries(summaries):
    notes = {}
    
    for header, summary in summaries.items():
        prompt = f"List the important information: {summary}"
        content_response = gemini_model.generate_content(prompt)
        
        content = content_response.text
        content = clean_content(content)
        
        notes[header] = content
    
    return notes

def clean_content(content):
    content = re.sub(r'\*+', '', content)  
    content = re.sub(r'\n+', '\n', content)  
    content = re.sub(r'\s+', ' ', content) 
    content = content.replace('-', '')
    bullet_points = content.strip().split('\n')
    bullet_points = [bp.strip() for bp in bullet_points if bp.strip()]
    return bullet_points


# Route to handle file upload and processing
@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' in request.files:
        file = request.files['file']
        if file.filename != '':
            uploads_dir = 'uploads'
            os.makedirs(uploads_dir, exist_ok=True)

            save_path = os.path.join(uploads_dir, file.filename)
            file.save(save_path)

            text = extract_text_from_presentation(save_path)
            sentences = extract_sentences_from_text(text)
            graph = create_graph(sentences)
            headers = create_headers(graph)
            summaries = summarize_clusters(headers)
            notes = generate_notes_from_summaries(summaries)

            return jsonify({
                'message': 'File uploaded and processed',
                'text': sentences,
                'graph_data': graph,
                'headers': headers,
                'summaries': summaries,
                'notes': notes,
                }), 200

    return jsonify({'error': 'No file part'}), 400

# Main function to run the app
if __name__ == '__main__':
    app.run(debug=True)