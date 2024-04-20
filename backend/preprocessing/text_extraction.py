from pptx import Presentation

# Function to extract text
def extract_text_from_presentation(ppt_file):
    prs = Presentation(ppt_file)
    texts = ''

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts += shape.text

    return texts